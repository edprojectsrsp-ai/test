"""Export Engine (Sprint 1) — PDF / DOCX / PPTX / XLSX from a shared payload.

Payload shape (friend-compatible + extensions)::

    {
      "title": str,
      "project_label": str,
      "fy_label": str,
      "month_label": str,
      "status_text": str,
      "header_lines": [str, ...],
      "physical_text": str,          # multiline
      "stage_text": str,
      "capex_text": str,
      "dpr_summary": [str, ...],
      "critical_rows": [[...], ...],  # activity tables
      "missed_rows": [[...], ...],
      "kpi_rows": [["Label", "Value"], ...],
      "table_sections": [
         {"title": str, "headers": [str,...], "rows": [[...], ...]},
         ...
      ],
      "current_view_image": optional path to PNG for screenshot export
    }
"""
from __future__ import annotations

import io
import os
import re
import tempfile
from datetime import datetime
from typing import Any, Dict, List, Optional, Sequence


# ── helpers ──────────────────────────────────────────────────────────────────

def _esc(s: Any) -> str:
    t = str(s if s is not None else "")
    return t.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")


def _lines(text: Any) -> List[str]:
    return [ln.strip() for ln in str(text or "").splitlines() if ln.strip()]


def _as_rows(rows: Any) -> List[List[str]]:
    out: List[List[str]] = []
    for r in rows or []:
        if isinstance(r, (list, tuple)):
            out.append([str(c) if c is not None else "" for c in r])
        else:
            out.append([str(r)])
    return out


# ── PDF ──────────────────────────────────────────────────────────────────────

def export_pdf(payload: Dict[str, Any], output_path: Optional[str] = None) -> bytes:
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4, landscape
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import (
        Image, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle,
    )

    buf = io.BytesIO()
    current_view_image = payload.get("current_view_image")
    page_size = landscape(A4) if current_view_image else A4
    doc = SimpleDocTemplate(
        buf if not output_path else output_path,
        pagesize=page_size,
        leftMargin=36, rightMargin=36, topMargin=36, bottomMargin=30,
    )
    styles = getSampleStyleSheet()
    story: list = []

    if current_view_image and os.path.isfile(str(current_view_image)):
        image = Image(str(current_view_image))
        scale = min(doc.width / max(image.imageWidth, 1), doc.height / max(image.imageHeight, 1))
        image.drawWidth = image.imageWidth * scale
        image.drawHeight = image.imageHeight * scale
        story.append(image)
        doc.build(story)
        if output_path:
            with open(output_path, "rb") as f:
                return f.read()
        return buf.getvalue()

    story.append(Paragraph(_esc(payload.get("title", "Executive Summary Dashboard")), styles["Title"]))
    story.append(Spacer(1, 6))
    for key, label in (
        ("project_label", "Project"),
        ("fy_label", "Financial Year"),
        ("month_label", "Month"),
        ("status_text", "Status"),
    ):
        story.append(Paragraph(f"<b>{label}:</b> {_esc(payload.get(key, '-'))}", styles["BodyText"]))
    story.append(Spacer(1, 10))

    def section_text(title: str, text: Any) -> None:
        story.append(Paragraph(f"<b>{_esc(title)}</b>", styles["Heading3"]))
        for line in _lines(text):
            story.append(Paragraph(_esc(line), styles["BodyText"]))
        story.append(Spacer(1, 8))

    if payload.get("header_lines"):
        story.append(Paragraph("<b>Project Identity</b>", styles["Heading3"]))
        for line in payload.get("header_lines") or []:
            if str(line).strip():
                story.append(Paragraph(_esc(line), styles["BodyText"]))
        story.append(Spacer(1, 8))

    if payload.get("kpi_rows"):
        story.append(Paragraph("<b>Key Indicators</b>", styles["Heading3"]))
        data = [["Metric", "Value"]] + _as_rows(payload["kpi_rows"])
        t = Table(data, repeatRows=1, colWidths=[220, 220])
        t.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#0b3d91")),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("GRID", (0, 0), (-1, -1), 0.4, colors.grey),
            ("FONTSIZE", (0, 0), (-1, -1), 9),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ]))
        story.append(t)
        story.append(Spacer(1, 10))

    if payload.get("physical_text"):
        section_text("Physical Progress Summary", payload.get("physical_text"))
    if payload.get("stage_text"):
        section_text("Stage Status", payload.get("stage_text"))
    if payload.get("capex_text"):
        section_text("CAPEX Snapshot", payload.get("capex_text"))

    if payload.get("dpr_summary"):
        story.append(Paragraph("<b>DPR Insights</b>", styles["Heading3"]))
        for item in payload.get("dpr_summary") or []:
            story.append(Paragraph(f"• {_esc(item)}", styles["BodyText"]))
        story.append(Spacer(1, 8))

    def add_table(title: str, headers: Sequence[str], rows: Sequence, header_color: str) -> None:
        if not rows:
            return
        story.append(Paragraph(f"<b>{_esc(title)}</b>", styles["Heading3"]))
        data = [list(headers)] + _as_rows(rows)
        t = Table(data, repeatRows=1)
        t.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor(header_color)),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("GRID", (0, 0), (-1, -1), 0.4, colors.grey),
            ("FONTSIZE", (0, 0), (-1, -1), 8),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ]))
        story.append(t)
        story.append(Spacer(1, 10))

    add_table(
        "Critical Path Activities",
        ["Activity", "Baseline Start", "Baseline Finish", "Current Status", "Delay"],
        payload.get("critical_rows") or [],
        "#0b3d91",
    )
    add_table(
        "Missed Baseline Activities",
        ["Activity", "Type", "Baseline Date", "Current Status"],
        payload.get("missed_rows") or [],
        "#7c2d12",
    )

    for sec in payload.get("table_sections") or []:
        headers = sec.get("headers") or []
        rows = sec.get("rows") or []
        if headers and rows:
            add_table(sec.get("title") or "Table", headers, rows, "#163f68")

    story.append(Spacer(1, 12))
    story.append(Paragraph(
        f"<font size='8' color='#666'>Generated {datetime.now().strftime('%d-%b-%Y %H:%M')} · Project Brain</font>",
        styles["BodyText"],
    ))

    doc.build(story)
    if output_path:
        with open(output_path, "rb") as f:
            return f.read()
    return buf.getvalue()


# ── DOCX ─────────────────────────────────────────────────────────────────────

def export_docx(payload: Dict[str, Any], output_path: Optional[str] = None) -> bytes:
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()
    current_view_image = payload.get("current_view_image")
    if current_view_image and os.path.isfile(str(current_view_image)):
        section = doc.sections[0]
        section.left_margin = Inches(0.35)
        section.right_margin = Inches(0.35)
        section.top_margin = Inches(0.35)
        section.bottom_margin = Inches(0.35)
        usable = section.page_width - section.left_margin - section.right_margin
        doc.add_picture(str(current_view_image), width=usable)
        if output_path:
            doc.save(output_path)
            with open(output_path, "rb") as f:
                return f.read()
        bio = io.BytesIO()
        doc.save(bio)
        return bio.getvalue()

    title = doc.add_heading(str(payload.get("title") or "Executive Summary Dashboard"), 0)
    title.alignment = WD_ALIGN_PARAGRAPH.LEFT

    for key, label in (
        ("project_label", "Project"),
        ("fy_label", "Financial Year"),
        ("month_label", "Month"),
        ("status_text", "Status"),
    ):
        doc.add_paragraph(f"{label}: {payload.get(key, '-')}")

    if payload.get("header_lines"):
        doc.add_heading("Project Identity", level=1)
        for line in payload.get("header_lines") or []:
            if str(line).strip():
                doc.add_paragraph(str(line))

    if payload.get("kpi_rows"):
        doc.add_heading("Key Indicators", level=1)
        table = doc.add_table(rows=1, cols=2)
        table.rows[0].cells[0].text = "Metric"
        table.rows[0].cells[1].text = "Value"
        for row in _as_rows(payload["kpi_rows"]):
            cells = table.add_row().cells
            cells[0].text = row[0] if row else ""
            cells[1].text = row[1] if len(row) > 1 else ""

    for heading, key in (
        ("Physical Progress Summary", "physical_text"),
        ("Stage Status", "stage_text"),
        ("CAPEX Snapshot", "capex_text"),
    ):
        text = payload.get(key)
        if text:
            doc.add_heading(heading, level=1)
            for line in _lines(text):
                doc.add_paragraph(line)

    if payload.get("dpr_summary"):
        doc.add_heading("DPR Insights", level=1)
        for item in payload.get("dpr_summary") or []:
            doc.add_paragraph(str(item), style="List Bullet")

    def add_table(title: str, headers: Sequence[str], rows: Sequence) -> None:
        if not rows:
            return
        doc.add_heading(title, level=1)
        table = doc.add_table(rows=1, cols=len(headers))
        for i, h in enumerate(headers):
            table.rows[0].cells[i].text = str(h)
        for row in _as_rows(rows):
            cells = table.add_row().cells
            for i, h in enumerate(headers):
                cells[i].text = row[i] if i < len(row) else ""

    add_table(
        "Critical Path Activities",
        ["Activity", "Baseline Start", "Baseline Finish", "Current Status", "Delay"],
        payload.get("critical_rows") or [],
    )
    add_table(
        "Missed Baseline Activities",
        ["Activity", "Type", "Baseline Date", "Current Status"],
        payload.get("missed_rows") or [],
    )
    for sec in payload.get("table_sections") or []:
        headers = sec.get("headers") or []
        rows = sec.get("rows") or []
        if headers and rows:
            add_table(sec.get("title") or "Table", headers, rows)

    foot = doc.add_paragraph(
        f"Generated {datetime.now().strftime('%d-%b-%Y %H:%M')} · Project Brain"
    )
    for run in foot.runs:
        run.font.size = Pt(8)
        run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    if output_path:
        doc.save(output_path)
        with open(output_path, "rb") as f:
            return f.read()
    bio = io.BytesIO()
    doc.save(bio)
    return bio.getvalue()


# ── PPTX ─────────────────────────────────────────────────────────────────────

def export_pptx(payload: Dict[str, Any], output_path: Optional[str] = None) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def add_title_bar(slide, title: str) -> None:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(13.333), Inches(0.8),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x0B, 0x3D, 0x91)
        shape.line.fill.background()
        box = slide.shapes.add_textbox(Inches(0.4), Inches(0.18), Inches(12), Inches(0.5))
        tf = box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.size = Pt(22)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    def add_bullets(slide, title: str, lines: List[str], left, top, width, height, size=12) -> None:
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = box.text_frame
        tf.word_wrap = True
        p0 = tf.paragraphs[0]
        r0 = p0.add_run()
        r0.text = title
        r0.font.size = Pt(16)
        r0.font.bold = True
        r0.font.color.rgb = RGBColor(0x0B, 0x3D, 0x91)
        items = lines or ["No data available."]
        for line in items[:18]:
            p = tf.add_paragraph()
            run = p.add_run()
            run.text = f"• {line}"
            run.font.size = Pt(size)

    def blank_slide():
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5),
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(0xEE, 0xF3, 0xF8)
        bg.line.fill.background()
        return slide

    slide = blank_slide()
    add_title_bar(slide, str(payload.get("title") or "Executive Summary Dashboard"))

    meta_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(6.2), Inches(1.5))
    tf = meta_box.text_frame
    meta_lines = [
        f"Project: {payload.get('project_label') or '-'}",
        f"Financial Year: {payload.get('fy_label') or '-'}",
        f"Month: {payload.get('month_label') or '-'}",
        f"Status: {payload.get('status_text') or '-'}",
    ]
    for i, line in enumerate(meta_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run() if not p.runs else p.runs[0]
        if i == 0 and p.runs:
            p.runs[0].text = line
            p.runs[0].font.size = Pt(14)
        else:
            run = p.add_run()
            run.text = line
            run.font.size = Pt(14)

    add_bullets(slide, "Project Identity",
                [str(x) for x in (payload.get("header_lines") or []) if str(x).strip()],
                0.5, 2.6, 5.7, 4.2, 12)
    add_bullets(slide, "Stage Status", _lines(payload.get("stage_text")), 6.5, 2.6, 3.0, 4.2, 11)
    add_bullets(slide, "CAPEX Snapshot", _lines(payload.get("capex_text")), 9.7, 2.6, 3.2, 4.2, 11)

    slide2 = blank_slide()
    add_bullets(slide2, "Physical Progress Summary", _lines(payload.get("physical_text")), 0.4, 0.3, 6.2, 4.5, 10)
    add_bullets(
        slide2, "DPR Insights",
        [str(x) for x in (payload.get("dpr_summary") or [])],
        6.9, 0.3, 5.9, 2.6, 12,
    )
    critical_lines = [
        f"{r[0]} | {r[3] if len(r) > 3 else ''} | Delay: {r[4] if len(r) > 4 else ''}"
        for r in _as_rows(payload.get("critical_rows") or [])[:12]
    ]
    add_bullets(slide2, "Critical Path Activities", critical_lines, 6.9, 3.1, 5.9, 3.8, 10)

    for sec in (payload.get("table_sections") or [])[:2]:
        slide_n = blank_slide()
        headers = sec.get("headers") or []
        rows = _as_rows(sec.get("rows") or [])
        lines = []
        if headers:
            lines.append(" | ".join(str(h) for h in headers))
        for r in rows[:20]:
            lines.append(" | ".join(r))
        add_bullets(slide_n, sec.get("title") or "Data", lines, 0.4, 0.4, 12.4, 6.6, 11)

    bio = io.BytesIO()
    if output_path:
        prs.save(output_path)
        with open(output_path, "rb") as f:
            return f.read()
    prs.save(bio)
    return bio.getvalue()


# ── XLSX ─────────────────────────────────────────────────────────────────────

def export_xlsx(payload: Dict[str, Any], output_path: Optional[str] = None) -> bytes:
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

    wb = Workbook()
    ws = wb.active
    ws.title = "Summary"

    header_fill = PatternFill("solid", fgColor="0B3D91")
    header_font = Font(color="FFFFFF", bold=True, size=11)
    thin = Border(
        left=Side(style="thin", color="B0B0B0"),
        right=Side(style="thin", color="B0B0B0"),
        top=Side(style="thin", color="B0B0B0"),
        bottom=Side(style="thin", color="B0B0B0"),
    )

    ws["A1"] = payload.get("title") or "Export"
    ws["A1"].font = Font(bold=True, size=14, color="0B3D91")
    ws.merge_cells("A1:D1")

    meta = [
        ("Project", payload.get("project_label")),
        ("Financial Year", payload.get("fy_label")),
        ("Month", payload.get("month_label")),
        ("Status", payload.get("status_text")),
    ]
    r = 3
    for label, val in meta:
        ws.cell(r, 1, label).font = Font(bold=True)
        ws.cell(r, 2, val or "-")
        r += 1

    r += 1
    if payload.get("kpi_rows"):
        ws.cell(r, 1, "Key Indicators").font = Font(bold=True, size=12)
        r += 1
        ws.cell(r, 1, "Metric").font = header_font
        ws.cell(r, 1).fill = header_fill
        ws.cell(r, 2, "Value").font = header_font
        ws.cell(r, 2).fill = header_fill
        r += 1
        for row in _as_rows(payload["kpi_rows"]):
            ws.cell(r, 1, row[0] if row else "")
            ws.cell(r, 2, row[1] if len(row) > 1 else "")
            r += 1
        r += 1

    def write_section(title: str, headers: Sequence[str], rows: Sequence) -> None:
        nonlocal r
        if not rows:
            return
        ws.cell(r, 1, title).font = Font(bold=True, size=12)
        r += 1
        for c, h in enumerate(headers, 1):
            cell = ws.cell(r, c, str(h))
            cell.font = header_font
            cell.fill = header_fill
            cell.border = thin
        r += 1
        for row in _as_rows(rows):
            for c, h in enumerate(headers, 1):
                cell = ws.cell(r, c, row[c - 1] if c - 1 < len(row) else "")
                cell.border = thin
                cell.alignment = Alignment(wrap_text=True, vertical="top")
            r += 1
        r += 1

    write_section(
        "Critical Path Activities",
        ["Activity", "Baseline Start", "Baseline Finish", "Current Status", "Delay"],
        payload.get("critical_rows") or [],
    )
    write_section(
        "Missed Baseline Activities",
        ["Activity", "Type", "Baseline Date", "Current Status"],
        payload.get("missed_rows") or [],
    )

    for i, sec in enumerate(payload.get("table_sections") or []):
        headers = sec.get("headers") or []
        rows = sec.get("rows") or []
        title = sec.get("title") or f"Table {i + 1}"
        # also put large tables on their own sheet
        if len(rows) > 15 and headers:
            safe = re.sub(r"[^\w\- ]", "", title)[:28] or f"Sheet{i+1}"
            wsn = wb.create_sheet(safe)
            for c, h in enumerate(headers, 1):
                cell = wsn.cell(1, c, str(h))
                cell.font = header_font
                cell.fill = header_fill
            for ri, row in enumerate(_as_rows(rows), 2):
                for c, h in enumerate(headers, 1):
                    wsn.cell(ri, c, row[c - 1] if c - 1 < len(row) else "")
            for col in wsn.columns:
                wsn.column_dimensions[col[0].column_letter].width = 16
        else:
            write_section(title, headers, rows)

    from openpyxl.utils import get_column_letter
    for idx in range(1, ws.max_column + 1):
        ws.column_dimensions[get_column_letter(idx)].width = 22

    if output_path:
        wb.save(output_path)
        with open(output_path, "rb") as f:
            return f.read()
    bio = io.BytesIO()
    wb.save(bio)
    return bio.getvalue()


# ── dispatcher ───────────────────────────────────────────────────────────────

RENDERERS = {
    "pdf": export_pdf,
    "docx": export_docx,
    "doc": export_docx,
    "pptx": export_pptx,
    "ppt": export_pptx,
    "xlsx": export_xlsx,
    "excel": export_xlsx,
}

MIME = {
    "pdf": "application/pdf",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "doc": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "ppt": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "excel": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
}

EXT = {
    "pdf": "pdf", "docx": "docx", "doc": "docx",
    "pptx": "pptx", "ppt": "pptx",
    "xlsx": "xlsx", "excel": "xlsx",
}


def render(payload: Dict[str, Any], fmt: str) -> tuple[bytes, str, str]:
    """Return (bytes, mime, filename_ext)."""
    key = (fmt or "pdf").strip().lower()
    if key not in RENDERERS:
        raise ValueError(f"Unsupported format: {fmt}")
    data = RENDERERS[key](payload)
    return data, MIME[key], EXT[key]


def safe_filename(stem: str, ext: str) -> str:
    clean = re.sub(r"[^\w\-]+", "_", stem or "export").strip("_")[:80] or "export"
    return f"{clean}.{ext}"
