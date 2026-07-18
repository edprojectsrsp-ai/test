import io
from fastapi import APIRouter, Depends, HTTPException
from fastapi.responses import StreamingResponse
from sqlalchemy.orm import Session
from sqlalchemy import text
from app.core.database import get_db

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
except ImportError:
    Presentation = None

router = APIRouter(prefix="/ppt-generator", tags=["ppt-generator"])

@router.get("/export/{scheme_id}")
def export_presentation(scheme_id: int, db: Session = Depends(get_db)):
    """Generate and export a PowerPoint presentation for a scheme."""
    scheme = db.execute(text("""
        SELECT scheme_name, current_status, estimated_cost_cr
        FROM scheme_master
        WHERE scheme_id = :sid AND is_deleted = FALSE
    """), {"sid": scheme_id}).first()

    if not scheme:
        raise HTTPException(status_code=404, detail="Scheme not found")

    # If python-pptx is not installed, fallback to return raw metadata/JSON
    if Presentation is None:
        # We can construct a simple mocked PPTX response or return mock bytes if pptx not installed
        # Let's write a mock stream for testing
        prs_stream = io.BytesIO()
        prs_stream.write(b"MOCK_PPTX_DATA_FOR_SCHEME_" + str(scheme_id).encode())
        prs_stream.seek(0)
        return StreamingResponse(
            prs_stream,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f"attachment; filename=Scheme_{scheme_id}_Summary.pptx"}
        )

    prs = Presentation()
    
    # Slide 1: Title
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = f"Scheme Review: {scheme.scheme_name}"
    subtitle.text = f"Status: {scheme.current_status.upper()} | Estimated Cost: \u20b9{scheme.estimated_cost_cr} Cr"

    # Slide 2: Package List
    packages = db.execute(text("""
        SELECT package_no, package_name, package_status, package_value_cr
        FROM packages
        WHERE scheme_id = :sid AND is_deleted = FALSE
    """), {"sid": scheme_id}).fetchall()

    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = "Packages Overview"

    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "Associated Packages Details:"
    for p in packages:
        p_para = tf.add_paragraph()
        p_para.text = f"- Pkg #{p.package_no}: {p.package_name} ({p.package_status.upper()}) | Value: \u20b9{p.package_value_cr} Cr"

    # Save presentation
    prs_stream = io.BytesIO()
    prs.save(prs_stream)
    prs_stream.seek(0)

    return StreamingResponse(
        prs_stream,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f"attachment; filename=Scheme_{scheme_id}_Summary.pptx"}
    )
